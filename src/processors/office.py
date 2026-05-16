"""Office docs (.docx, .pptx, .xlsx, .odt, .epub, ...) via MarkItDown with native fallbacks."""
from pathlib import Path
from typing import Any, Dict


def _via_markitdown(path: Path) -> str:
    from markitdown import MarkItDown

    md = MarkItDown()
    result = md.convert(str(path))
    return getattr(result, "text_content", "") or getattr(result, "markdown", "")


def _docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"=== Slide {i} ===")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text)
    return "\n".join(parts)


def _xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"=== Sheet: {ws.title} ===")
        for row in ws.iter_rows(values_only=True):
            parts.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(parts)


def process(path: Path) -> Dict[str, Any]:
    ext = path.suffix.lower()
    text = ""

    try:
        text = _via_markitdown(path)
    except Exception:
        text = ""

    if not text.strip():
        try:
            if ext == ".docx":
                text = _docx(path)
            elif ext == ".pptx":
                text = _pptx(path)
            elif ext == ".xlsx":
                text = _xlsx(path)
        except Exception as e:
            text = f"[office extraction failed: {e}]"

    return {
        "kind": "office",
        "source": str(path),
        "filename": path.name,
        "text": text,
    }
